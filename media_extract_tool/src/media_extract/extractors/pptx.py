from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from importlib.metadata import PackageNotFoundError, version
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from media_extract.udr import new_block_id, now_iso


def _pptx_version() -> str:
    try:
        return version("python-pptx")
    except PackageNotFoundError:
        return "unknown"


def _shape_text(shape: Any) -> str:
    if shape.has_text_frame:
        return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        tbl = shape.table
        lines: list[str] = []
        for row in tbl.rows:
            lines.append(" | ".join(c.text.strip() for c in row.cells))
        return "\n".join(lines).strip()
    return ""


def extract_pptx(path: Path, source_uri: str) -> dict[str, Any]:
    raw = path.read_bytes()
    fp = hashlib.sha256(raw).hexdigest()
    document: dict[str, Any] = {
        "sourceFormat": "pptx",
        "sourceUri": source_uri,
        "fingerprint": fp,
        "originalFilename": path.name,
        "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "byteSize": len(raw),
        "ingestedAt": now_iso(),
        "pipeline": {
            "libraryVersions": {"python-pptx": _pptx_version()},
        },
        "pptx": {"slideCount": 0},
    }

    prs = Presentation(str(path))
    slides = list(prs.slides)
    document["pptx"]["slideCount"] = len(slides)

    blocks: list[dict[str, Any]] = []
    n = 0

    def push(btype: str, text: str, slide_index: int, shape_id: str | None = None) -> None:
        nonlocal n
        if not text.strip():
            return
        n += 1
        loc: dict[str, Any] = {"slideIndex": slide_index}
        if shape_id:
            loc["shapeId"] = shape_id
        blocks.append(
            {
                "blockId": new_block_id(n),
                "type": btype,
                "text": text.strip(),
                "confidence": None,
                "location": loc,
            }
        )

    for si, slide in enumerate(slides):
        for shape_idx, shape in enumerate(slide.shapes):
            sid = f"s{si}_shape{shape_idx}"
            txt = _shape_text(shape)
            if not txt:
                continue
            if shape_idx == 0 and shape.has_text_frame:
                push("slide_title", txt, si, sid)
            else:
                push("slide_body", txt, si, sid)
        try:
            notes = slide.notes_slide
            if notes and notes.notes_text_frame:
                nt = notes.notes_text_frame.text.strip()
                if nt:
                    n += 1
                    blocks.append(
                        {
                            "blockId": new_block_id(n),
                            "type": "speaker_notes",
                            "text": nt,
                            "confidence": None,
                            "location": {"slideIndex": si, "shapeId": "notes"},
                        }
                    )
        except Exception:
            pass

    return {"schemaVersion": "1.0.0", "document": document, "blocks": blocks, "assets": []}
